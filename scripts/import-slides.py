#!/usr/bin/env python3
"""
One-time import of the source deck into the app.

Why this exists
---------------
Every stage cited "deck slides 21-35", and nineteen figure captions cited
"Slide 23 - ...", but the deck shipped nowhere. Those were dead pointers: they
told a reader to go and look at something they did not have. This script makes
them real by bringing the slides into the project as web images plus text.

What it produces
----------------
  public/slides/slide-NN.webp   98 rendered slides
  src/slides.ts                 generated metadata: number, title, text, stage

The text matters as much as the picture. It gives every slide a truthful alt
attribute rather than "Slide 23", it puts the deck's own words into the app's
search index, and it means the slide is still useful to somebody who cannot see
the image or is reading with the images stripped out.

Prerequisites (import only - the outputs are committed, so a normal build and
a normal checkout never need any of this):

    apt-get install libreoffice poppler-utils
    pip install python-pptx pillow

Usage:

    python3 scripts/import-slides.py path/to/deck.pptx

Rendering path: LibreOffice to PDF, poppler to PNG, Pillow to WebP. Going via
PDF rather than exporting images straight from LibreOffice keeps vector text
crisp - it is rasterised once, at the resolution we ask for, instead of being
resampled from whatever LibreOffice felt like emitting.
"""

import json
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation

# 1280px wide renders 16:9 slides comfortably above the ~1100px the lightbox
# ever displays, so there is headroom to zoom without them going soft. Quality
# 74 is where the WebP curve flattens for slide content - flat fills and text,
# rather than photographs. Together: ~26KB a slide.
WIDTH = 1280
QUALITY = 74

ROOT = Path(__file__).resolve().parent.parent


def stage_ranges() -> list[tuple[str, int, int]]:
    """
    Read the module ids and their slide ranges straight out of course.ts.

    An earlier version of this script hardcoded the list, and two of the nine
    ids were guesses that did not exist ("stakeholders" and "maturity" for what
    are really "government" and "integration"). Thirteen slides ended up filed
    against nothing: their stage filter showed an empty grid. Reading the one
    place the mapping actually lives makes that class of mistake impossible.
    """
    source = (ROOT / "src" / "course.ts").read_text(encoding="utf8")
    pattern = re.compile(
        r'id:\s*"([a-z]+)",\s*\n\s*number:\s*\d+,[\s\S]{0,400}?slides:\s*"(\d+)\s*[–—-]\s*(\d+)"'
    )
    found = [(m.group(1), int(m.group(2)), int(m.group(3))) for m in pattern.finditer(source)]
    if len(found) != 9:
        raise SystemExit(f"Expected 9 stage ranges in course.ts, parsed {len(found)}")
    return found


RANGES = stage_ranges()


def stage_for(number: int) -> str:
    for stage, first, last in RANGES:
        if first <= number <= last:
            return stage
    raise ValueError(f"Slide {number} falls outside every stage range")


def clean(text: str) -> str:
    """Collapse whitespace and drop the slide-number text box."""
    text = text.replace("\v", " ").replace("\xa0", " ")
    text = re.sub(r"\s+", " ", text).strip()
    return text


def slide_strings(slide) -> list[str]:
    """Every readable string on a slide, in shape order, tables included."""
    out: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            value = clean(shape.text_frame.text)
            if value:
                out.append(value)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                cells = [clean(cell.text) for cell in row.cells]
                cells = [cell for cell in cells if cell]
                # Tables repeat merged cells across the span; de-duplicate so a
                # merged header does not appear four times in the search index.
                deduped: list[str] = []
                for cell in cells:
                    if not deduped or deduped[-1] != cell:
                        deduped.append(cell)
                if deduped:
                    out.append(" | ".join(deduped))
    return out


def render(pptx: Path, out_dir: Path, expected: int) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)
    for old in out_dir.glob("slide-*.webp"):
        old.unlink()

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        shutil.copy(pptx, tmp_path / "deck.pptx")
        subprocess.run(
            ["soffice", "--headless", "--norestore", "--convert-to", "pdf",
             "--outdir", str(tmp_path), str(tmp_path / "deck.pptx")],
            check=True, capture_output=True,
            # LibreOffice writes a profile; point it somewhere disposable so it
            # never touches the invoking user's real one.
            env={"HOME": str(tmp_path), "PATH": "/usr/bin:/bin"},
        )
        pdf = tmp_path / "deck.pdf"
        if not pdf.exists():
            raise SystemExit("LibreOffice produced no PDF.")

        # 108 dpi against a 960pt-wide page lands just over our target width;
        # Pillow does the final, better-quality resize down to exactly WIDTH.
        subprocess.run(
            ["pdftoppm", "-r", "108", "-png", str(pdf), str(tmp_path / "page")],
            check=True, capture_output=True,
        )

        pages = sorted(tmp_path.glob("page-*.png"))
        if len(pages) != expected:
            raise SystemExit(f"Rendered {len(pages)} pages but the deck has {expected} slides.")

        for page in pages:
            number = int(page.stem.split("-")[1])
            image = Image.open(page).convert("RGB")
            if image.width != WIDTH:
                height = round(image.height * WIDTH / image.width)
                image = image.resize((WIDTH, height), Image.LANCZOS)
            image.save(out_dir / f"slide-{number:02d}.webp", "WEBP",
                       quality=QUALITY, method=6)


def main() -> None:
    if len(sys.argv) != 2:
        raise SystemExit("Usage: python3 scripts/import-slides.py <deck.pptx>")
    pptx = Path(sys.argv[1]).expanduser()
    if not pptx.exists():
        raise SystemExit(f"No such file: {pptx}")

    presentation = Presentation(str(pptx))
    slides = list(presentation.slides)

    records = []
    for number, slide in enumerate(slides, 1):
        strings = slide_strings(slide)
        # The deck stamps its own slide number in a text box on most slides.
        # It is noise once the number is structural, so drop it.
        strings = [s for s in strings if s != str(number)]

        # Prefer the real title placeholder. Taking the first shape instead is
        # wrong whenever a designer put a body box above the heading in z-order
        # - slide 2 came out titled "Discover user needs and opportunities..."
        # when the heading was "Session Outcomes".
        title = ""
        try:
            placeholder = slide.shapes.title
            if placeholder is not None:
                title = clean(placeholder.text)
        except (AttributeError, KeyError):
            pass
        # Some slides carry no words at all — slide 98 is a wordless "questions?"
        # illustration. Leave the title empty rather than inventing one or
        # falling back to "Slide 98", which would give the useless alt text
        # "Slide 98: Slide 98". The UI says "image only" instead.
        if not title or title == str(number):
            title = strings[0] if strings else ""

        body = [s for s in strings if s != title]
        records.append({
            "n": number,
            "stage": stage_for(number),
            "title": title,
            "text": " ".join(body),
        })

    out_dir = ROOT / "public" / "slides"
    render(pptx, out_dir, len(slides))

    total = sum(f.stat().st_size for f in out_dir.glob("*.webp"))

    lines = [
        "/**",
        " * Generated by scripts/import-slides.py - do not edit by hand.",
        " *",
        " * Metadata for the source deck. The images live in public/slides and are",
        " * attached at build time; this file carries the words, which is what makes",
        " * a slide searchable, describable to a screen reader, and readable when the",
        " * picture has not loaded.",
        " */",
        "",
        "export type Slide = {",
        "  /** 1-based slide number, as cited throughout the course. */",
        "  n: number;",
        "  /** Module id of the stage this slide belongs to. */",
        "  stage: string;",
        "  /** The slide's own heading. Empty when the slide carries no words. */",
        "  title: string;",
        "  /** Everything else on the slide, flattened. */",
        "  text: string;",
        "};",
        "",
        f"export const SLIDE_COUNT = {len(records)};",
        "",
        "export const slides: Slide[] = [",
    ]
    for record in records:
        lines.append("  {")
        lines.append(f"    n: {record['n']},")
        lines.append(f"    stage: {json.dumps(record['stage'])},")
        lines.append(f"    title: {json.dumps(record['title'])},")
        lines.append(f"    text: {json.dumps(record['text'])},")
        lines.append("  },")
    lines.append("];")
    lines.append("")

    (ROOT / "src" / "slides.ts").write_text("\n".join(lines), encoding="utf8")

    print(f"  {len(records)} slides -> public/slides ({total / 1024 / 1024:.2f} MB, "
          f"{total / len(records) / 1024:.0f} KB mean)")
    print("  src/slides.ts written")


if __name__ == "__main__":
    main()
